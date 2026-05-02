"""
Google Slides build pipeline.

Reads a unit's stage JSONs + composed PNGs and produces a Google Slides deck
that the PM can review/edit.

First-run requires OAuth interaction (browser will open on your machine to
grant Drive + Slides scopes). Subsequent runs use the cached token.

Usage:
    from pipeline.slides import build_unit_deck
    url = build_unit_deck(Path("generated_units/batch_1/k_patterns_pattern_parade"))
    print(url)
"""

from __future__ import annotations
import json
import os
from pathlib import Path
from typing import Any

from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

PROJECT_ROOT = Path(__file__).parent.parent
CRED_PATH = PROJECT_ROOT / "credentials.json"
TOKEN_PATH = PROJECT_ROOT / "token.json"

SCOPES = [
    "https://www.googleapis.com/auth/drive",
    "https://www.googleapis.com/auth/presentations",
]

# Parent Drive folder where every newly-built unit deck lands. Each unit
# gets its own subfolder named after the blueprint's thematic_title; the
# deck and any uploaded composite images for that unit go inside.
# Override per-call with ``build_unit_deck(..., drive_parent_folder_id=...)``.
# Set 2026-04-29 to the TCE "Google Slides Version" shared folder.
UNIT_DECK_PARENT_FOLDER_ID = os.environ.get(
    "TCE_UNIT_DECK_PARENT_FOLDER_ID",
    "1hKMJcWOZDvsksHPSA7jPkpsd_6E4aSQ2",
)

# Slide dimensions — PORTRAIT 7.5in × 10in (matches the Little Programmers PPTX).
SLIDE_WIDTH_EMU = 6858000   # 7.5 in
SLIDE_HEIGHT_EMU = 9144000  # 10 in
SLIDE_WIDTH_IN = 7.5
SLIDE_HEIGHT_IN = 10.0

EMU_PER_IN = 914400
EMU_PER_PT = 12700

# ── Typography (matches Little Programmers PPTX exactly) ──────────────────
# Confirmed by inspecting the PPTX XML:
#   - Worksheet titles use "Chelsea Market" (the playful handwritten font)
#   - Body, lesson titles, and most other text use "Lexend"
#   - Cover subtitle uses "Lexend ExtraLight"
TITLE_FONT = "Chelsea Market"      # kid-facing playful titles (worksheets, certificate)
HEADING_FONT = "Lexend"            # cover, overview, lesson plan, manipulative, rubric, marketplace titles
BODY_FONT = "Lexend"               # body text everywhere
SUBTITLE_FONT = "Lexend"           # cover subtitle (with ExtraLight weight via API not directly supported, fallback to gray)
GRAY_LABEL = {"red": 0.45, "green": 0.45, "blue": 0.45}
BLACK = {"red": 0, "green": 0, "blue": 0}


def emu(inches: float) -> int:
    return int(inches * EMU_PER_IN)


# ── Auth ─────────────────────────────────────────────────────────────────

def get_credentials() -> Credentials:
    """Load cached token or trigger OAuth flow."""
    creds = None
    if TOKEN_PATH.exists():
        creds = Credentials.from_authorized_user_file(str(TOKEN_PATH), SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            print("Refreshing expired token...")
            creds.refresh(Request())
        else:
            print("First-run: opening browser for OAuth grant. Approve Drive + Slides scopes...")
            flow = InstalledAppFlow.from_client_secrets_file(str(CRED_PATH), SCOPES)
            creds = flow.run_local_server(port=0)
        TOKEN_PATH.write_text(creds.to_json())
        print(f"Token saved to {TOKEN_PATH}")
    return creds


# ── Drive image upload (cached) ───────────────────────────────────────────

class DriveUploader:
    def __init__(self, drive_service, folder_id: str | None = None):
        self.drive = drive_service
        self.folder_id = folder_id
        self._cache: dict[str, str] = {}  # path -> Drive file ID

    def upload(self, image_path: Path, name: str | None = None) -> str:
        """Upload a PNG, make it readable by anyone with the link, return file ID."""
        key = str(image_path.resolve())
        if key in self._cache:
            return self._cache[key]
        media = MediaFileUpload(str(image_path), mimetype="image/png")
        body: dict = {"name": name or image_path.name}
        if self.folder_id:
            body["parents"] = [self.folder_id]
        f = self.drive.files().create(body=body, media_body=media, fields="id").execute()
        file_id = f["id"]
        # Make it readable so Slides API can fetch it
        self.drive.permissions().create(
            fileId=file_id,
            body={"role": "reader", "type": "anyone"},
        ).execute()
        self._cache[key] = file_id
        return file_id

    def public_url(self, file_id: str) -> str:
        return f"https://drive.google.com/uc?id={file_id}&export=download"


def _find_or_create_subfolder(drive, parent_id: str, name: str) -> str:
    """Find a Drive folder named ``name`` under ``parent_id``, creating it if
    absent. Returns the folder's file ID. Trailing whitespace in ``name`` is
    stripped before matching to avoid duplicates.
    """
    name = name.strip()
    # Escape single quotes for the Drive query language.
    safe_name = name.replace("'", "\\'")
    q = (
        f"'{parent_id}' in parents and "
        f"mimeType = 'application/vnd.google-apps.folder' and "
        f"name = '{safe_name}' and trashed = false"
    )
    res = drive.files().list(q=q, fields="files(id,name)", pageSize=10).execute()
    for f in res.get("files", []):
        return f["id"]
    created = drive.files().create(
        body={
            "name": name,
            "mimeType": "application/vnd.google-apps.folder",
            "parents": [parent_id],
        },
        fields="id",
    ).execute()
    return created["id"]


# ── Slide-building helpers (return batchUpdate request lists) ─────────────

def _gen_id(prefix: str, slide_idx: int, elem_idx: int = 0) -> str:
    """Stable element IDs make later edits easier."""
    return f"{prefix}_{slide_idx:03d}_{elem_idx:02d}"


def _create_slide_request(slide_id: str, layout: str = "BLANK") -> dict:
    return {
        "createSlide": {
            "objectId": slide_id,
            "slideLayoutReference": {"predefinedLayout": layout},
        }
    }


def _create_text_box(slide_id: str, elem_id: str, x_in: float, y_in: float,
                     w_in: float, h_in: float, text: str,
                     font_size_pt: float = 14, bold: bool = False,
                     align: str = "START", font_family: str = BODY_FONT,
                     color: dict | None = None) -> list[dict]:
    style: dict = {
        "fontFamily": font_family,
        "fontSize": {"magnitude": font_size_pt, "unit": "PT"},
        "bold": bold,
    }
    fields = "fontFamily,fontSize,bold"
    if color is not None:
        style["foregroundColor"] = {"opaqueColor": {"rgbColor": color}}
        fields += ",foregroundColor"

    requests = [
        {
            "createShape": {
                "objectId": elem_id,
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {
                        "width": {"magnitude": emu(w_in), "unit": "EMU"},
                        "height": {"magnitude": emu(h_in), "unit": "EMU"},
                    },
                    "transform": {
                        "scaleX": 1, "scaleY": 1,
                        "translateX": emu(x_in), "translateY": emu(y_in),
                        "unit": "EMU",
                    },
                },
            }
        },
        {"insertText": {"objectId": elem_id, "insertionIndex": 0, "text": text}},
        {
            "updateTextStyle": {
                "objectId": elem_id,
                "style": style,
                "fields": fields,
                "textRange": {"type": "ALL"},
            }
        },
        {
            "updateParagraphStyle": {
                "objectId": elem_id,
                "style": {"alignment": align},
                "fields": "alignment",
                "textRange": {"type": "ALL"},
            }
        },
    ]
    return requests


def _create_boxed_banner(slide_id: str, elem_id: str, x_in: float, y_in: float,
                         w_in: float, h_in: float, text: str,
                         font_size_pt: float = 13,
                         bold_prefix: str | None = None,
                         dashed: bool = False) -> list[dict]:
    """Create a rectangle shape with a thin black border and text inside.
    Mirrors the PDF's "Learning Goal" boxed banner.
    If bold_prefix is given, only that prefix is bolded (e.g., 'Learning Goal:')."""
    requests = [
        {
            "createShape": {
                "objectId": elem_id,
                "shapeType": "RECTANGLE",
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {
                        "width": {"magnitude": emu(w_in), "unit": "EMU"},
                        "height": {"magnitude": emu(h_in), "unit": "EMU"},
                    },
                    "transform": {
                        "scaleX": 1, "scaleY": 1,
                        "translateX": emu(x_in), "translateY": emu(y_in),
                        "unit": "EMU",
                    },
                },
            }
        },
        # Outline + transparent fill
        {
            "updateShapeProperties": {
                "objectId": elem_id,
                "shapeProperties": {
                    "shapeBackgroundFill": {
                        "solidFill": {"color": {"rgbColor": {"red": 1, "green": 1, "blue": 1}}, "alpha": 1}
                    },
                    "outline": {
                        "weight": {"magnitude": (1.5 if not dashed else 1), "unit": "PT"},
                        "outlineFill": {"solidFill": {"color": {"rgbColor": BLACK}}},
                        "dashStyle": "DASH" if dashed else "SOLID",
                    },
                },
                "fields": "shapeBackgroundFill,outline",
            }
        },
    ]
    # Only insert + style text if there's text to add (decorative borders pass empty text)
    if text:
        requests += [
            {"insertText": {"objectId": elem_id, "insertionIndex": 0, "text": text}},
            {
                "updateTextStyle": {
                    "objectId": elem_id,
                    "style": {
                        "fontFamily": BODY_FONT,
                        "fontSize": {"magnitude": font_size_pt, "unit": "PT"},
                        "bold": False,
                    },
                    "fields": "fontFamily,fontSize,bold",
                    "textRange": {"type": "ALL"},
                }
            },
            {
                "updateParagraphStyle": {
                    "objectId": elem_id,
                    "style": {"alignment": "START"},
                    "fields": "alignment",
                    "textRange": {"type": "ALL"},
                }
            },
        ]
        # Bold the prefix if given
        if bold_prefix and bold_prefix in text:
            end = len(bold_prefix)
            requests.append({
                "updateTextStyle": {
                    "objectId": elem_id,
                    "style": {"bold": True},
                    "fields": "bold",
                    "textRange": {"type": "FIXED_RANGE", "startIndex": 0, "endIndex": end},
                }
            })
    return requests


def _create_underline(slide_id: str, elem_id: str, x_in: float, y_in: float,
                      w_in: float, weight_pt: float = 1.0) -> list[dict]:
    """Thin horizontal line, used as the writing line beside Name: / Date: labels."""
    return [
        {
            "createLine": {
                "objectId": elem_id,
                "lineCategory": "STRAIGHT",
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {
                        "width": {"magnitude": emu(w_in), "unit": "EMU"},
                        "height": {"magnitude": 0, "unit": "EMU"},
                    },
                    "transform": {
                        "scaleX": 1, "scaleY": 1,
                        "translateX": emu(x_in), "translateY": emu(y_in),
                        "unit": "EMU",
                    },
                },
            }
        },
        {
            "updateLineProperties": {
                "objectId": elem_id,
                "lineProperties": {
                    "lineFill": {"solidFill": {"color": {"rgbColor": GRAY_LABEL}}},
                    "weight": {"magnitude": weight_pt, "unit": "PT"},
                },
                "fields": "lineFill,weight",
            }
        },
    ]


def _create_name_date_row(slide_id: str, idx: int, y_in: float = 0.40) -> list[dict]:
    """Standard worksheet header: 'Name: _____   Date: _____' in gray.
    Sized for portrait 7.5in slide width."""
    requests: list[dict] = []
    requests += _create_text_box(
        slide_id, _gen_id("nm_lbl", idx, 0),
        x_in=0.40, y_in=y_in, w_in=1.0, h_in=0.32,
        text="Name:", font_size_pt=14, font_family=BODY_FONT, color=GRAY_LABEL,
    )
    requests += _create_underline(
        slide_id, _gen_id("nm_ln", idx, 0),
        x_in=1.20, y_in=y_in + 0.30, w_in=2.55,
    )
    requests += _create_text_box(
        slide_id, _gen_id("dt_lbl", idx, 0),
        x_in=4.10, y_in=y_in, w_in=1.0, h_in=0.32,
        text="Date:", font_size_pt=14, font_family=BODY_FONT, color=GRAY_LABEL,
    )
    requests += _create_underline(
        slide_id, _gen_id("dt_ln", idx, 0),
        x_in=4.90, y_in=y_in + 0.30, w_in=2.30,
    )
    return requests


def _create_image(slide_id: str, elem_id: str, file_id: str,
                  x_in: float, y_in: float, w_in: float, h_in: float) -> dict:
    return {
        "createImage": {
            "objectId": elem_id,
            "url": f"https://drive.google.com/uc?id={file_id}&export=download",
            "elementProperties": {
                "pageObjectId": slide_id,
                "size": {
                    "width": {"magnitude": emu(w_in), "unit": "EMU"},
                    "height": {"magnitude": emu(h_in), "unit": "EMU"},
                },
                "transform": {
                    "scaleX": 1, "scaleY": 1,
                    "translateX": emu(x_in), "translateY": emu(y_in),
                    "unit": "EMU",
                },
            },
        }
    }


def _create_table(slide_id: str, elem_id: str, x_in: float, y_in: float,
                  w_in: float, h_in: float, rows: int, cols: int) -> dict:
    return {
        "createTable": {
            "objectId": elem_id,
            "elementProperties": {
                "pageObjectId": slide_id,
                "size": {
                    "width": {"magnitude": emu(w_in), "unit": "EMU"},
                    "height": {"magnitude": emu(h_in), "unit": "EMU"},
                },
                "transform": {
                    "scaleX": 1, "scaleY": 1,
                    "translateX": emu(x_in), "translateY": emu(y_in),
                    "unit": "EMU",
                },
            },
            "rows": rows,
            "columns": cols,
        }
    }


def _insert_table_text(table_id: str, row: int, col: int, text: str,
                       font_size_pt: float = 11, bold: bool = False) -> list[dict]:
    return [
        {
            "insertText": {
                "objectId": table_id,
                "cellLocation": {"rowIndex": row, "columnIndex": col},
                "text": text,
                "insertionIndex": 0,
            }
        },
        {
            "updateTextStyle": {
                "objectId": table_id,
                "cellLocation": {"rowIndex": row, "columnIndex": col},
                "style": {
                    "fontFamily": "Arial",
                    "fontSize": {"magnitude": font_size_pt, "unit": "PT"},
                    "bold": bold,
                },
                "fields": "fontFamily,fontSize,bold",
                "textRange": {"type": "ALL"},
            }
        },
    ]


# ── Layout helpers (auto-fit / overflow-safe sizing) ──────────────────────

def _est_lines(text: str, max_chars_per_line: int) -> int:
    """Rough line count for a text block in a box of given character width.
    Counts paragraph wraps + newline-induced lines. Used to size boxes that
    must hold variable-length text without overflow."""
    if not text:
        return 0
    total = 0
    for p in text.split("\n"):
        if not p.strip():
            total += 1
            continue
        total += max(1, (len(p) + max_chars_per_line - 1) // max_chars_per_line)
    return total


def _fit_font_size(text: str, base_pt: float, w_in: float, h_in: float,
                   min_pt: float | None = None) -> float:
    """Shrink font size so that `text` fits within (w_in × h_in) at `base_pt`.
    Heuristic: chars/line ≈ w_in * 165 / pt; line height ≈ pt * 1.4 / 72.
    Returns a font size between min_pt and base_pt."""
    if min_pt is None:
        min_pt = base_pt * 0.55
    pt = base_pt
    while pt > min_pt:
        chars_per_line = max(8, int(w_in * 165 / pt))
        n_lines = _est_lines(text, chars_per_line)
        line_h_in = pt * 1.4 / 72
        if n_lines * line_h_in <= h_in:
            return pt
        pt -= 1
    return min_pt


# ── Per-slide-type builders ───────────────────────────────────────────────

def build_cover_slide(idx: int, slide_id: str, bp: dict, mk: dict,
                      uploader: DriveUploader, composed_dir: Path) -> list[dict]:
    """Cover slide — portrait 7.5×10. Matches the Little Programmers PPTX:
    big boxed title with stacked text upper-middle, smaller boxed subtitle in gray below,
    character at bottom centered.

    Title pulls the grade from `bp['grade']` (was previously hardcoded to
    'Kindergarten' — fixed 2026-04-28 to support G1/G2/G3). The subtitle box
    auto-shrinks the font for thematic titles longer than ~30 characters so
    long titles like 'The Pattern Parade — Pattern Rules and Number Patterns
    to 50' do not overflow the box border."""
    requests = [_create_slide_request(slide_id)]

    grade_text = bp["grade"].strip()
    big_box_text = f"{grade_text}\nOntario Curriculum\n{bp['strand'].split('. ')[-1]}"
    # Auto-fit the big-box font (grade name varies in length).
    BIG_BOX_W, BIG_BOX_H = 6.4, 2.45
    BIG_FONT_SZ = _fit_font_size(big_box_text, base_pt=40, w_in=BIG_BOX_W,
                                 h_in=BIG_BOX_H - 0.30, min_pt=26)
    # Big boxed title — centered horizontally, upper-third
    requests += _create_boxed_banner(
        slide_id, _gen_id("big_box", idx, 0),
        x_in=0.55, y_in=2.55, w_in=BIG_BOX_W, h_in=BIG_BOX_H,
        text=big_box_text, font_size_pt=BIG_FONT_SZ,
    )
    requests.append({
        "updateTextStyle": {
            "objectId": _gen_id("big_box", idx, 0),
            "style": {"fontFamily": HEADING_FONT, "fontSize": {"magnitude": BIG_FONT_SZ, "unit": "PT"}},
            "fields": "fontFamily,fontSize",
            "textRange": {"type": "ALL"},
        }
    })
    requests.append({
        "updateParagraphStyle": {
            "objectId": _gen_id("big_box", idx, 0),
            "style": {"alignment": "CENTER"},
            "fields": "alignment",
            "textRange": {"type": "ALL"},
        }
    })

    # Smaller boxed subtitle — Lexend, gray, centered. Adaptive to thematic title length.
    subtitle = bp["thematic_title"]
    SUB_W = 5.6
    if len(subtitle) > 50:
        SUB_H, SM_FONT_SZ = 1.45, 16
    elif len(subtitle) > 30:
        SUB_H, SM_FONT_SZ = 1.20, 18
    else:
        SUB_H, SM_FONT_SZ = 0.85, 24
    SUB_X = (SLIDE_WIDTH_IN - SUB_W) / 2
    requests += _create_boxed_banner(
        slide_id, _gen_id("sm_box", idx, 0),
        x_in=SUB_X, y_in=5.30, w_in=SUB_W, h_in=SUB_H,
        text=subtitle, font_size_pt=SM_FONT_SZ,
    )
    requests.append({
        "updateTextStyle": {
            "objectId": _gen_id("sm_box", idx, 0),
            "style": {"fontFamily": HEADING_FONT,
                      "fontSize": {"magnitude": SM_FONT_SZ, "unit": "PT"},
                      "foregroundColor": {"opaqueColor": {"rgbColor": GRAY_LABEL}}},
            "fields": "fontFamily,fontSize,foregroundColor",
            "textRange": {"type": "ALL"},
        }
    })
    requests.append({
        "updateParagraphStyle": {
            "objectId": _gen_id("sm_box", idx, 0),
            "style": {"alignment": "CENTER"},
            "fields": "alignment",
            "textRange": {"type": "ALL"},
        }
    })

    # Coco at bottom-centre — placed below the (possibly enlarged) subtitle box
    coco_path = composed_dir / "CHAR_COCO_FRONT.png"
    if coco_path.exists():
        coco_id = uploader.upload(coco_path, "char_coco_front.png")
        tw, th = 1.6, 2.1
        ix = (SLIDE_WIDTH_IN - tw) / 2
        coco_y = max(6.80, 5.30 + SUB_H + 0.30)
        requests.append(_create_image(slide_id, _gen_id("img", idx, 0), coco_id,
                                      x_in=ix, y_in=coco_y, w_in=tw, h_in=th))

    return requests


def build_overview_slide(idx: int, slide_id: str, bp: dict) -> list[dict]:
    """Unit overview — portrait. Matches PPTX slide 2 (Lexend 23pt title, body 13pt).

    Title height now grows with the thematic title length so long titles like
    'Unit Overview: The Pattern Parade — Pattern Rules and Number Patterns to
    50' do not overlap the Grade/Strand line below (fixed 2026-04-28). All
    subsequent body lines reposition based on the actual title height."""
    requests = [_create_slide_request(slide_id)]
    title_text = f"Unit Overview: {bp['thematic_title']}"
    TITLE_W = 6.8
    title_pt = _fit_font_size(title_text, base_pt=23, w_in=TITLE_W, h_in=2.10, min_pt=18)
    chars_per_line = max(8, int(TITLE_W * 165 / title_pt))
    n_title_lines = max(1, _est_lines(title_text, chars_per_line))
    title_h = max(0.65, n_title_lines * (title_pt * 1.4 / 72) + 0.10)
    requests += _create_text_box(
        slide_id, _gen_id("ttl", idx, 0),
        x_in=0.35, y_in=0.30, w_in=TITLE_W, h_in=title_h,
        text=title_text,
        font_size_pt=title_pt, bold=True, font_family=HEADING_FONT,
    )
    # Body block — bold labels with bold prefixes. Stack starts after the title.
    cur_y = 0.30 + title_h + 0.20
    requests += _create_text_box(
        slide_id, _gen_id("gs", idx, 0),
        x_in=0.35, y_in=cur_y, w_in=6.8, h_in=0.35,
        text=f"Grade: {bp['grade']}  |  Strand: {bp['strand'].split('. ')[-1]}",
        font_size_pt=13, font_family=BODY_FONT,
    )
    # Bold the "Grade:" and "Strand:" labels
    requests.append({
        "updateTextStyle": {
            "objectId": _gen_id("gs", idx, 0),
            "style": {"bold": True}, "fields": "bold",
            "textRange": {"type": "FIXED_RANGE", "startIndex": 0, "endIndex": 6},
        }
    })
    cur_y += 0.40
    requests += _create_text_box(
        slide_id, _gen_id("te", idx, 0),
        x_in=0.35, y_in=cur_y, w_in=6.8, h_in=0.35,
        text=f"Target Expectations: {', '.join(bp['curriculum_codes'])}",
        font_size_pt=13, font_family=BODY_FONT,
    )
    requests.append({
        "updateTextStyle": {
            "objectId": _gen_id("te", idx, 0),
            "style": {"bold": True}, "fields": "bold",
            "textRange": {"type": "FIXED_RANGE", "startIndex": 0, "endIndex": 20},
        }
    })
    cur_y += 0.45
    lg_text = f"Learning Goal: {bp['unit_learning_goal']}"
    lg_lines = _est_lines(lg_text, max_chars_per_line=80)
    lg_h = max(0.55, lg_lines * 0.22)
    requests += _create_text_box(
        slide_id, _gen_id("lg", idx, 0),
        x_in=0.35, y_in=cur_y, w_in=6.8, h_in=lg_h,
        text=lg_text,
        font_size_pt=13, font_family=BODY_FONT,
    )
    requests.append({
        "updateTextStyle": {
            "objectId": _gen_id("lg", idx, 0),
            "style": {"bold": True}, "fields": "bold",
            "textRange": {"type": "FIXED_RANGE", "startIndex": 0, "endIndex": 14},
        }
    })
    table_y = cur_y + lg_h + 0.25
    rows = len(bp["lesson_arc"]) + 1
    cols = 3  # PPTX has 3 cols: Lesson | Lesson Title | Learning Goal
    table_id = _gen_id("tbl", idx, 0)
    table_h = max(3.0, SLIDE_HEIGHT_IN - table_y - 0.45)
    requests.append(_create_table(slide_id, table_id, x_in=0.35, y_in=table_y,
                                  w_in=6.8, h_in=table_h, rows=rows, cols=cols))
    return requests, table_id


# ── Lesson-split dispatch (DISABLED per user preference, 2026-05-02) ─────
#
# The split-layout helpers below (`build_lesson_slide_split`,
# `build_lesson_slides`) are kept in this module so the option can be
# re-enabled by editing `lesson_should_split`. They are NOT currently
# used: `lesson_should_split` returns False unconditionally, so every
# lesson renders on a single slide using `build_lesson_slide`. The
# proportional-scaling code inside `build_lesson_slide` already absorbs
# G3 density adequately for the single-slide format.
#
# If a future grade or content shape ever overflows again, change
# `lesson_should_split` to return True for that case — the dispatcher
# wiring in `build_unit_deck` already supports two-slide lessons.

GRADES_REQUIRING_SPLIT_LAYOUT: set[str] = set()  # empty = never auto-split by grade


def lesson_should_split(lp: dict, grade: str | None = None) -> bool:
    """Decide whether this lesson renders on one slide or splits across two.

    Currently disabled by user preference: always returns False, so every
    lesson uses the single-slide layout regardless of grade or step count.
    Toggle by editing this function.
    """
    return False


def _action_steps_summary(lp: dict, *, slice_start: int = 0,
                           slice_end: int | None = None) -> str:
    """Action body text built from a slice of step instructions.

    Used by the split-layout to emit the first half of the steps on slide A
    and the rest on slide B. Same first-sentence/180-char treatment as
    `_lesson_section_summary` for consistency."""
    s = lp["action"]
    steps = s["steps"]
    if slice_end is None:
        slice_end = len(steps)
    chunk = steps[slice_start:slice_end]
    bullets = [f"• {_first_sentence(st['instruction'], max_chars=200)}" for st in chunk]
    if slice_start == 0:
        return f"{s['activity_name']}\n\n" + "\n".join(bullets)
    return "\n".join(bullets)


def _lesson_section_summary(lp: dict, key: str) -> str:
    """Teacher-facing summary for a lesson section — paragraph-length per the PPTX style.

    All caps lifted on 2026-04-28: action shows EVERY step (was capped at 3) and
    consolidation shows EVERY discussion prompt (was capped at 2). The lesson
    slide layout dynamically sizes each section's body height based on this
    text via _est_lines, so longer Grade 2/3 content no longer leaves a
    dead-air gap between Action and Consolidation."""
    s = lp[key]
    if key == "minds_on":
        # Hook + a paragraph excerpt from the teacher script. Excerpt grows
        # for higher-grade scripts (G2 ~500c, G3 ~600c) up to ~360 chars; the
        # full text remains in the JSON / unit.md.
        script = s.get("teacher_script", "")
        para = script.split("\n\n")[0] if "\n\n" in script else script
        if len(para) > 360:
            para = para[:360].rsplit(" ", 1)[0] + "…"
        return f"{s['activity_name']}\n\n{s['hook']}\n\n{para}"
    if key == "action":
        # ALL step bullets (no cap). First sentence per step, max 180 chars.
        steps = s["steps"]
        bullets = [f"• {_first_sentence(st['instruction'], max_chars=180)}" for st in steps]
        return f"{s['activity_name']}\n\n" + "\n".join(bullets)
    if key == "consolidation":
        # ALL discussion prompts + exit routine (no cap)
        prompts = [_first_sentence(p, max_chars=130) for p in s.get("discussion_prompts", [])]
        prompt_lines = "\n".join(f"  • {p}" for p in prompts)
        exit_routine = _first_sentence(s.get("exit_routine", ""), max_chars=200)
        return f"{s['activity_name']}\n\nDiscuss:\n{prompt_lines}\n\nExit: {exit_routine}"
    return ""


def _first_sentence(text: str, max_chars: int = 140) -> str:
    """Return the first sentence of a string, truncated to max_chars."""
    if not text:
        return ""
    text = text.strip()
    # Split on sentence enders, take first
    for end in (". ", "! ", "? "):
        if end in text:
            text = text.split(end, 1)[0] + end.strip()
            break
    if len(text) > max_chars:
        text = text[:max_chars - 1].rsplit(" ", 1)[0] + "…"
    return text


def _truncate_to_height(text: str, h_in: float,
                         line_in: float = 0.20, w_chars: int = 90) -> str:
    """Truncate ``text`` so the rendered block fits within ``h_in`` inches at
    the given line height. Counts wrapped lines per paragraph using the same
    estimator as the layout. Used by lesson sections when the proportional
    scaler shrinks a body height below the text's natural height — without
    truncation the overflow runs off the bottom of the slide.

    Preserves full lines (no mid-word cuts on the wrap boundary). If the
    truncation lands inside a paragraph, the paragraph keeps its leading
    sentences and ends with an ellipsis on the last retained line.
    """
    if not text or h_in <= 0:
        return text or ""
    max_lines = max(1, int(h_in / line_in))
    if max_lines <= 0:
        return ""

    out_lines: list[str] = []
    used_lines = 0
    paragraphs = text.split("\n")
    for p in paragraphs:
        if used_lines >= max_lines:
            break
        # If empty line, count it once (matches _est_lines behaviour).
        if not p.strip():
            out_lines.append("")
            used_lines += 1
            continue
        n_lines = _est_lines(p, w_chars)
        if used_lines + n_lines <= max_lines:
            out_lines.append(p)
            used_lines += n_lines
            continue
        # Need to truncate this paragraph. Keep up to the line budget,
        # ending with an ellipsis on the final retained line.
        budget = max_lines - used_lines
        # Soft estimate: budget * w_chars characters can fit
        char_budget = max(20, budget * w_chars - 1)
        truncated = p[:char_budget].rsplit(" ", 1)[0]
        if not truncated:
            truncated = p[:char_budget]
        out_lines.append(truncated.rstrip(",.;: ").strip() + "…")
        used_lines = max_lines
        break
    return "\n".join(out_lines)


def build_lesson_slide(idx: int, slide_id: str, lp: dict) -> list[dict]:
    """Lesson plan slide — portrait. Matches PPTX slide 3 exactly:
    Lexend 20pt title, then bold-labelled metadata lines, then numbered sections."""
    n = lp["lesson_number"]
    requests = [_create_slide_request(slide_id)]

    # Title — Lexend 20pt bold
    requests += _create_text_box(
        slide_id, _gen_id("ttl", idx, 0),
        x_in=0.35, y_in=0.30, w_in=6.8, h_in=0.55,
        text=f"Lesson {n}: {lp['lesson_title']}",
        font_size_pt=20, bold=True, font_family=HEADING_FONT,
    )
    # Bold-prefix metadata lines (Goal:, Expectation:, Learning Goal:, Materials:)
    y = 1.00
    line_h = 0.45
    metadata_lines = [
        ("Goal: ", lp["big_idea"]),
        ("Expectation: ", ", ".join(lp["primary_expectations"])),
        ("Learning Goal: ", f"\"{lp['student_learning_goal']}\""),
        ("Materials: ", ", ".join(lp["manipulatives_used"][:4])),
    ]
    for i, (label, body) in enumerate(metadata_lines):
        elem_id = _gen_id(f"meta{i}", idx, 0)
        full = f"{label}{body}"
        requests += _create_text_box(
            slide_id, elem_id,
            x_in=0.35, y_in=y + i * line_h, w_in=6.8, h_in=line_h,
            text=full, font_size_pt=13, font_family=BODY_FONT,
        )
        # Bold the label prefix
        requests.append({
            "updateTextStyle": {
                "objectId": elem_id,
                "style": {"bold": True}, "fields": "bold",
                "textRange": {"type": "FIXED_RANGE", "startIndex": 0, "endIndex": len(label)},
            }
        })

    # Three numbered sections, single-column prose. Heights are computed from
    # the actual content (lifted from a fixed 2.20"-per-section budget on
    # 2026-04-28 so Action with fewer steps no longer leaves a dead gap before
    # Consolidation, and longer G2/G3 content can grow without overflowing).
    sections = [
        ("minds_on", "1. Minds On"),
        ("action", "2. Action"),
        ("consolidation", "3. Consolidation"),
    ]
    section_y = y + len(metadata_lines) * line_h + 0.25  # ~3.05
    available = SLIDE_HEIGHT_IN - section_y - 0.30        # leave 0.30" bottom margin
    HEADER_H = 0.40
    GAP_H = 0.18
    BODY_LINE_IN = 0.20  # ~12pt at line-height 1.4
    BODY_W_CHARS = 90    # at 12pt across 6.6" usable width

    # Pre-compute body text + estimated heights
    body_metrics = []
    for key, label in sections:
        s = lp[key]
        body = _lesson_section_summary(lp, key)
        if s["activity_name"] in body:
            body = body.replace(s["activity_name"], "", 1).lstrip("\n")
        body = body.strip()
        n_lines = max(2, _est_lines(body, BODY_W_CHARS))
        natural_h = max(0.50, n_lines * BODY_LINE_IN)
        body_metrics.append({"key": key, "label": label, "body": body,
                             "duration": s["duration_minutes"], "activity": s["activity_name"],
                             "natural_h": natural_h})

    # Total natural height
    total_natural = sum(HEADER_H + m["natural_h"] + GAP_H for m in body_metrics)
    if total_natural > available:
        # Scale body heights proportionally so everything fits
        slack_needed = total_natural - available
        scalable = sum(m["natural_h"] for m in body_metrics)
        for m in body_metrics:
            shrink = (m["natural_h"] / scalable) * slack_needed
            m["body_h"] = max(0.50, m["natural_h"] - shrink)
            # CRITICAL: also truncate the body text to fit the shrunken box.
            # Without this, the text rendering overflows the box bounds and
            # bleeds into the next section's space (or off the slide edge).
            # Keeps a small safety margin (BODY_LINE_IN) so the last visible
            # line isn't kissing the bottom edge of the box.
            m["body"] = _truncate_to_height(
                m["body"],
                h_in=max(0.40, m["body_h"] - BODY_LINE_IN * 0.5),
                line_in=BODY_LINE_IN,
                w_chars=BODY_W_CHARS,
            )
    else:
        # Distribute remaining space proportionally so sections fill the slide
        # rather than ending with white space at the bottom.
        extra = available - total_natural
        natural_total = sum(m["natural_h"] for m in body_metrics)
        for m in body_metrics:
            m["body_h"] = m["natural_h"] + extra * (m["natural_h"] / natural_total)

    # Place sections sequentially
    sy = section_y
    for i, m in enumerate(body_metrics):
        header_text = f"{m['label']} ({m['duration']} mins): {m['activity']}"
        requests += _create_text_box(
            slide_id, _gen_id(f"shdr{i}", idx, 0),
            x_in=0.35, y_in=sy, w_in=6.8, h_in=HEADER_H,
            text=header_text, font_size_pt=14, bold=True, font_family=HEADING_FONT,
        )
        requests += _create_text_box(
            slide_id, _gen_id(f"sbod{i}", idx, 0),
            x_in=0.55, y_in=sy + HEADER_H + 0.02, w_in=6.6, h_in=m["body_h"],
            text=m["body"], font_size_pt=12, font_family=BODY_FONT,
        )
        sy += HEADER_H + m["body_h"] + GAP_H
    return requests


def build_lesson_slide_split(idx_a: int, idx_b: int,
                              slide_id_a: str, slide_id_b: str,
                              lp: dict) -> list[dict]:
    """Two-slide lesson layout for dense (G3+) content.

    Slide A — header + metadata + Minds On + Action header + first half of
    Action steps. Slide B — continuation header + remaining Action steps +
    Consolidation. Same fonts and proportional scaling as the single-slide
    layout, but each slide gets its own SLIDE_HEIGHT_IN of vertical budget.
    """
    n = lp["lesson_number"]
    title = lp["lesson_title"]
    requests: list[dict] = []
    steps = lp["action"]["steps"]
    half = (len(steps) + 1) // 2  # round up; e.g. 6 → 3+3, 7 → 4+3
    a_steps = steps[:half]
    b_steps = steps[half:]

    # ── SLIDE A ──
    requests.append(_create_slide_request(slide_id_a))
    requests += _create_text_box(
        slide_id_a, _gen_id("ttl_a", idx_a, 0),
        x_in=0.35, y_in=0.30, w_in=6.8, h_in=0.55,
        text=f"Lesson {n}: {title} — Part 1",
        font_size_pt=20, bold=True, font_family=HEADING_FONT,
    )
    y = 1.00
    line_h = 0.45
    metadata_lines = [
        ("Goal: ", lp["big_idea"]),
        ("Expectation: ", ", ".join(lp["primary_expectations"])),
        ("Learning Goal: ", f"\"{lp['student_learning_goal']}\""),
        ("Materials: ", ", ".join(lp["manipulatives_used"][:4])),
    ]
    for i, (label, body) in enumerate(metadata_lines):
        elem_id = _gen_id(f"meta_a{i}", idx_a, 0)
        full = f"{label}{body}"
        requests += _create_text_box(
            slide_id_a, elem_id,
            x_in=0.35, y_in=y + i * line_h, w_in=6.8, h_in=line_h,
            text=full, font_size_pt=13, font_family=BODY_FONT,
        )
        requests.append({
            "updateTextStyle": {
                "objectId": elem_id,
                "style": {"bold": True}, "fields": "bold",
                "textRange": {"type": "FIXED_RANGE", "startIndex": 0, "endIndex": len(label)},
            }
        })

    # Two sections on slide A: Minds On (full) + Action (first half)
    sa_y = y + len(metadata_lines) * line_h + 0.25
    available_a = SLIDE_HEIGHT_IN - sa_y - 0.30
    HEADER_H = 0.40
    GAP_H = 0.18
    BODY_LINE_IN = 0.20
    BODY_W_CHARS = 90

    # Minds On body
    mo = lp["minds_on"]
    mo_body = _lesson_section_summary(lp, "minds_on")
    if mo["activity_name"] in mo_body:
        mo_body = mo_body.replace(mo["activity_name"], "", 1).lstrip("\n")
    mo_body = mo_body.strip()

    # Action (first half) body
    act = lp["action"]
    a_body = _action_steps_summary(lp, slice_start=0, slice_end=half)
    if act["activity_name"] in a_body:
        a_body = a_body.replace(act["activity_name"], "", 1).lstrip("\n")
    a_body = a_body.strip()

    sections_a = [
        {"label": "1. Minds On", "duration": mo["duration_minutes"],
         "activity": mo["activity_name"], "body": mo_body},
        {"label": f"2. Action — Part 1 (steps 1–{half})",
         "duration": act["duration_minutes"], "activity": act["activity_name"],
         "body": a_body},
    ]
    for m in sections_a:
        m["natural_h"] = max(0.50, max(2, _est_lines(m["body"], BODY_W_CHARS)) * BODY_LINE_IN)
    total_natural = sum(HEADER_H + m["natural_h"] + GAP_H for m in sections_a)
    if total_natural > available_a:
        slack = total_natural - available_a
        natural_total = sum(m["natural_h"] for m in sections_a)
        for m in sections_a:
            m["body_h"] = max(0.50, m["natural_h"] - slack * (m["natural_h"] / natural_total))
    else:
        extra = available_a - total_natural
        natural_total = sum(m["natural_h"] for m in sections_a)
        for m in sections_a:
            m["body_h"] = m["natural_h"] + extra * (m["natural_h"] / natural_total)

    sy = sa_y
    for i, m in enumerate(sections_a):
        header_text = f"{m['label']} ({m['duration']} mins): {m['activity']}"
        requests += _create_text_box(
            slide_id_a, _gen_id(f"shdr_a{i}", idx_a, 0),
            x_in=0.35, y_in=sy, w_in=6.8, h_in=HEADER_H,
            text=header_text, font_size_pt=14, bold=True, font_family=HEADING_FONT,
        )
        requests += _create_text_box(
            slide_id_a, _gen_id(f"sbod_a{i}", idx_a, 0),
            x_in=0.55, y_in=sy + HEADER_H + 0.02, w_in=6.6, h_in=m["body_h"],
            text=m["body"], font_size_pt=12, font_family=BODY_FONT,
        )
        sy += HEADER_H + m["body_h"] + GAP_H

    # Continued-on-next-slide footer
    requests += _create_text_box(
        slide_id_a, _gen_id("foot_a", idx_a, 0),
        x_in=0.35, y_in=SLIDE_HEIGHT_IN - 0.30, w_in=6.8, h_in=0.25,
        text=f"Lesson {n} continues on the next slide →",
        font_size_pt=10, font_family=BODY_FONT, color=GRAY_LABEL,
        align="END",
    )

    # ── SLIDE B ──
    requests.append(_create_slide_request(slide_id_b))
    requests += _create_text_box(
        slide_id_b, _gen_id("ttl_b", idx_b, 0),
        x_in=0.35, y_in=0.30, w_in=6.8, h_in=0.55,
        text=f"Lesson {n}: {title} — Part 2 (continued)",
        font_size_pt=20, bold=True, font_family=HEADING_FONT,
    )

    sb_y = 1.05
    available_b = SLIDE_HEIGHT_IN - sb_y - 0.30

    # Action (second half) body
    a_body_b = _action_steps_summary(lp, slice_start=half, slice_end=len(steps))
    a_body_b = a_body_b.strip()

    # Consolidation
    cs = lp["consolidation"]
    cs_body = _lesson_section_summary(lp, "consolidation")
    if cs["activity_name"] in cs_body:
        cs_body = cs_body.replace(cs["activity_name"], "", 1).lstrip("\n")
    cs_body = cs_body.strip()

    sections_b = [
        {"label": f"2. Action — Part 2 (steps {half + 1}–{len(steps)})",
         "duration": act["duration_minutes"], "activity": act["activity_name"],
         "body": a_body_b},
        {"label": "3. Consolidation", "duration": cs["duration_minutes"],
         "activity": cs["activity_name"], "body": cs_body},
    ]
    for m in sections_b:
        m["natural_h"] = max(0.50, max(2, _est_lines(m["body"], BODY_W_CHARS)) * BODY_LINE_IN)
    total_natural = sum(HEADER_H + m["natural_h"] + GAP_H for m in sections_b)
    if total_natural > available_b:
        slack = total_natural - available_b
        natural_total = sum(m["natural_h"] for m in sections_b)
        for m in sections_b:
            m["body_h"] = max(0.50, m["natural_h"] - slack * (m["natural_h"] / natural_total))
    else:
        extra = available_b - total_natural
        natural_total = sum(m["natural_h"] for m in sections_b)
        for m in sections_b:
            m["body_h"] = m["natural_h"] + extra * (m["natural_h"] / natural_total)

    sy = sb_y
    for i, m in enumerate(sections_b):
        header_text = f"{m['label']} ({m['duration']} mins): {m['activity']}"
        requests += _create_text_box(
            slide_id_b, _gen_id(f"shdr_b{i}", idx_b, 0),
            x_in=0.35, y_in=sy, w_in=6.8, h_in=HEADER_H,
            text=header_text, font_size_pt=14, bold=True, font_family=HEADING_FONT,
        )
        requests += _create_text_box(
            slide_id_b, _gen_id(f"sbod_b{i}", idx_b, 0),
            x_in=0.55, y_in=sy + HEADER_H + 0.02, w_in=6.6, h_in=m["body_h"],
            text=m["body"], font_size_pt=12, font_family=BODY_FONT,
        )
        sy += HEADER_H + m["body_h"] + GAP_H

    return requests


def build_lesson_slides(idx: int, slide_id_base: str, lp: dict,
                         grade: str | None = None) -> tuple[list[dict], int]:
    """Dispatcher: returns (request_list, slides_emitted_count).

    For G3+ (or any lesson with ≥6 action steps) this emits two slides;
    otherwise one. Caller increments its slide index by the second value.
    """
    if lesson_should_split(lp, grade):
        sid_a = f"{slide_id_base}_a"
        sid_b = f"{slide_id_base}_b"
        reqs = build_lesson_slide_split(idx, idx + 1, sid_a, sid_b, lp)
        return reqs, 2
    return build_lesson_slide(idx, slide_id_base, lp), 1


def _fit_image(png_path: Path, max_w_in: float, max_h_in: float) -> tuple[float, float]:
    """Return (target_w_in, target_h_in) that fits within the box AND preserves aspect ratio."""
    from PIL import Image as PILImage
    with PILImage.open(png_path) as pim:
        pw, ph = pim.size
    aspect = pw / ph  # >1 = wide
    # Try fitting to the max width first
    target_w = max_w_in
    target_h = target_w / aspect
    if target_h > max_h_in:
        target_h = max_h_in
        target_w = target_h * aspect
    return (target_w, target_h)


def build_worksheet_slide(idx: int, slide_id: str, ws: dict,
                          uploader: DriveUploader, composed_dir: Path) -> list[dict]:
    """Worksheet slide — portrait 7.5×10. Matches the Little Programmers PPTX slide 4:

      0.40  Name: ____         Date: ____             (Lexend 14pt gray)
      0.95  Centered Chelsea Market title              (28pt)
      1.85  Boxed Learning Goal banner                 (Lexend 16pt + bold prefix)
      2.95  "Instructions for the Student:" header     (Lexend 14pt bold)
      3.30  Numbered instructions (1./2./3.)           (Lexend 14pt with bold key terms)
      4.95  Hero image — centered                      (large, fills remaining space)
      9.55  Footer "(N parts on the printable...)"     (Lexend 10pt gray)
    """
    n = ws["lesson_number"]
    requests = [_create_slide_request(slide_id)]

    # 1. Name / Date row — gray Lexend with line-shape underlines
    requests += _create_name_date_row(slide_id, idx, y_in=0.40)

    # 2. Centered Chelsea Market title — large, like the PPTX. Auto-fit
    # so the title ALWAYS fits on one line. Chelsea Market is a much wider
    # display font than the default _fit_font_size heuristic assumes — even
    # at w_in=5.5 a 30-char title would still wrap at 30pt. We feed the
    # fitter an effective width of 4.5" AND drop the base to 28pt so the
    # calc stays in the right neighbourhood for Chelsea Market.
    # K title 'Welcome to the Parade!' (23 chars) → ~28pt, single line.
    # G3 'Lead the Number Parade to 1000' (30 chars) → ~24pt, single line.
    title_text = ws["worksheet_title"]
    title_pt = _fit_font_size(title_text, base_pt=28, w_in=4.5, h_in=0.85,
                               min_pt=16)
    requests += _create_text_box(
        slide_id, _gen_id("ttl", idx, 0),
        x_in=0.35, y_in=1.00, w_in=6.8, h_in=0.85,
        text=title_text,
        font_size_pt=title_pt, font_family=TITLE_FONT, align="CENTER",
    )

    # 3. Boxed "Learning Goal:" banner — wide, prominent
    goal_text = f"Learning Goal: {ws['student_learning_goal']}"
    requests += _create_boxed_banner(
        slide_id, _gen_id("goal_box", idx, 0),
        x_in=0.45, y_in=2.05, w_in=6.6, h_in=0.95,
        text=goal_text, font_size_pt=16, bold_prefix="Learning Goal:",
    )

    # 4. "Instructions for the Student:" section header
    parts = ws["pages"][0]["parts"]
    section_label = "Mission" if any("Mission" in p.get("part_title", "")
                                     or "Find" in p.get("part_title", "")
                                     or "Help" in p.get("part_title", "") for p in parts) else "Instructions for the Student:"
    requests += _create_text_box(
        slide_id, _gen_id("sec_hdr", idx, 0),
        x_in=0.35, y_in=3.20, w_in=6.8, h_in=0.40,
        text=section_label, font_size_pt=15, bold=True, font_family=HEADING_FONT,
    )

    # 5. Numbered instructions (1./2./3.)
    instr_lines = []
    for i, part in enumerate(parts[:3], 1):
        instr = _first_sentence(part["student_instructions"], max_chars=120)
        instr_lines.append(f"{i}. {instr}")
    instr_text = "\n".join(instr_lines)
    requests += _create_text_box(
        slide_id, _gen_id("instr", idx, 0),
        x_in=0.35, y_in=3.65, w_in=6.8, h_in=1.65,
        text=instr_text, font_size_pt=14, font_family=BODY_FONT,
    )

    # 6. Hero image — large, centered, fills the bottom half.
    # The picker now ranks ALL image_placeholders by composed-PNG file size
    # so that placeholder PNGs (compose.py emits ~2-3 KB grey rectangles when
    # it can't render real content) get out-ranked by genuine content. ID
    # match is still preferred but only among files >= 2.5 KB; below that
    # threshold the picker falls back to the largest PNG in the worksheet.
    PLACEHOLDER_BYTES = 2500

    candidates: list[dict] = []
    for part in parts:
        for ph in part.get("image_placeholders", []) or []:
            png = composed_dir / f"{ph['id']}.png"
            size = png.stat().st_size if png.exists() else 0
            id_match_score = 1 if any(
                k in ph["id"] for k in ("PARADE", "TOP", "BLANK", "ANIMAL", "BOX")
            ) else 0
            looks_real = size >= PLACEHOLDER_BYTES
            candidates.append({
                "ph": ph, "png_path": png, "size": size,
                "id_match": id_match_score, "looks_real": looks_real,
            })

    # Pick: prefer (real-looking AND id-match), then real-looking, then
    # any with the biggest file. Never pick a non-existent file.
    candidates = [c for c in candidates if c["png_path"].exists()]
    hero_ph = None
    if candidates:
        best = max(candidates, key=lambda c: (
            c["looks_real"], c["id_match"], c["size"]
        ))
        hero_ph = best["ph"]

    if hero_ph:
        png_path = composed_dir / f"{hero_ph['id']}.png"
        if png_path.exists():
            file_id = uploader.upload(png_path, png_path.name)
            # Image area: 0.35 to 7.15 wide (6.8) × 5.50 to 9.40 tall (3.90)
            tw, th = _fit_image(png_path, max_w_in=6.8, max_h_in=3.90)
            ix = 0.35 + (6.8 - tw) / 2
            iy = 5.50 + (3.90 - th) / 2
            requests.append(_create_image(slide_id, _gen_id("hero", idx, 0),
                                          file_id, x_in=ix, y_in=iy, w_in=tw, h_in=th))

    # 7. Parts footer
    if len(parts) > 1:
        footer_text = f"({len(parts)} parts on the printable worksheet)"
        requests += _create_text_box(
            slide_id, _gen_id("foot", idx, 0),
            x_in=0.35, y_in=9.55, w_in=6.8, h_in=0.30,
            text=footer_text, font_size_pt=10, font_family=BODY_FONT,
            color=GRAY_LABEL, align="CENTER",
        )

    return requests


def build_manipulative_slide(idx: int, slide_id: str, asset: dict,
                             uploader: DriveUploader, composed_dir: Path) -> list[dict]:
    """Manipulative slide — portrait. Centered title + image, 2-column specs/prep below."""
    requests = [_create_slide_request(slide_id)]
    # Centered title
    requests += _create_text_box(
        slide_id, _gen_id("ttl", idx, 0),
        x_in=0.35, y_in=0.30, w_in=6.8, h_in=0.55,
        text=f"{asset['name']}",
        font_size_pt=24, bold=True, font_family=HEADING_FONT, align="CENTER",
    )
    # Centered asset ID (gray)
    requests += _create_text_box(
        slide_id, _gen_id("aid", idx, 0),
        x_in=0.35, y_in=0.88, w_in=6.8, h_in=0.30,
        text=f"({asset['asset_id']})",
        font_size_pt=12, font_family=BODY_FONT, color=GRAY_LABEL, align="CENTER",
    )
    # Centered purpose
    purpose_short = _first_sentence(asset["purpose"], max_chars=200)
    requests += _create_text_box(
        slide_id, _gen_id("pur", idx, 0),
        x_in=0.7, y_in=1.20, w_in=6.1, h_in=0.75,
        text=purpose_short, font_size_pt=13, font_family=BODY_FONT, align="CENTER",
    )
    # Image — large, centered horizontally (image placement unchanged per user)
    if asset.get("image_placeholders"):
        ph = asset["image_placeholders"][0]
        png_path = composed_dir / f"{ph['id']}.png"
        if png_path.exists():
            file_id = uploader.upload(png_path, png_path.name)
            tw, th = _fit_image(png_path, max_w_in=6.8, max_h_in=3.50)
            ix = 0.35 + (6.8 - tw) / 2
            requests.append(_create_image(slide_id, _gen_id("img", idx, 0), file_id,
                                          x_in=ix, y_in=2.10, w_in=tw, h_in=th))

    # 2-column specs/prep below the image (uses the white space)
    # Left column: PRINT SPECS + USE IN CLASS
    # Right column: TEACHER PREP
    ps = asset["print_specifications"]
    qty = asset["quantity_per_class"]
    for delim in [" = ", " (", " — ", "; "]:
        if delim in qty:
            qty = qty.split(delim, 1)[0].strip()
            break
    qty = qty[:55].rstrip(",.;: ").strip()
    left_text = (
        f"PRINT SPECS\n"
        f"   Size: {ps['page_size']} {ps['orientation']}\n"
        f"   Colour: {ps['color']}\n"
        f"   Pages per set: {ps['pages_per_set']}\n"
        f"   Laminate: {'yes' if ps['laminate_recommended'] else 'no'}\n"
        f"\n"
        f"USE IN CLASS\n"
        f"   Quantity: {qty}\n"
        f"   Lessons: {', '.join(map(str, asset['used_in_lessons']))}\n"
        f"   Prep: ~{asset['estimated_prep_minutes']} min"
    )
    requests += _create_text_box(
        slide_id, _gen_id("spec", idx, 0),
        x_in=0.45, y_in=5.95, w_in=3.30, h_in=3.70,
        text=left_text, font_size_pt=12, font_family=BODY_FONT,
    )
    # TEACHER PREP — keep first sentences short enough that 5 steps fit
    # within the 3.70" tall right-column box. Each step is capped at 90
    # chars (was 110 — caused overflow on G3 manipulatives where prep
    # instructions reference Grade 3 specifics like "Lessons 3, 4, and 5
    # all require dry-erase number writing"). The whole block is then
    # truncated to box height as a final safety net.
    prep_steps = asset["teacher_prep_steps"][:5]
    prep = "TEACHER PREP\n" + "\n".join(
        f"   {i+1}. {_first_sentence(s, 90)}" for i, s in enumerate(prep_steps)
    )
    if len(asset["teacher_prep_steps"]) > 5:
        prep += "\n   …"
    PREP_H = 3.70
    prep = _truncate_to_height(prep, PREP_H - 0.10, line_in=0.20, w_chars=42)
    requests += _create_text_box(
        slide_id, _gen_id("prep", idx, 0),
        x_in=3.95, y_in=5.95, w_in=3.30, h_in=PREP_H,
        text=prep, font_size_pt=12, font_family=BODY_FONT,
    )
    return requests


def build_rubric_slide(idx: int, slide_id: str, rub: dict, bp: dict) -> list[dict]:
    """Rubric modelled on Little Programmers page 21:
    Name/Date row, centered playful title, then the 4×5 table."""
    requests = [_create_slide_request(slide_id)]
    requests += _create_name_date_row(slide_id, idx, y_in=0.40)
    # Centered Chelsea Market title
    requests += _create_text_box(
        slide_id, _gen_id("ttl", idx, 0),
        x_in=0.35, y_in=1.00, w_in=6.8, h_in=0.85,
        text="Summative Assessment Rubric", font_size_pt=30,
        font_family=TITLE_FONT, align="CENTER",
    )
    n_rows = len(rub["rows"]) + 1
    n_cols = 5
    table_id = _gen_id("rub_tbl", idx, 0)
    # Lots of vertical room in portrait
    requests.append(_create_table(slide_id, table_id, x_in=0.30, y_in=2.05,
                                  w_in=6.9, h_in=7.50, rows=n_rows, cols=n_cols))
    return requests, table_id


def build_certificate_slide(idx: int, slide_id: str, cert: dict, bp: dict,
                            uploader: DriveUploader, composed_dir: Path) -> list[dict]:
    """Certificate — portrait. Dashed outer border, big Chelsea Market title, centered prose."""
    requests = [_create_slide_request(slide_id)]

    # Dashed outer border (full slide minus margin)
    requests += _create_boxed_banner(
        slide_id, _gen_id("dashed", idx, 0),
        x_in=0.25, y_in=0.25, w_in=7.0, h_in=9.50,
        text="", font_size_pt=8, dashed=True,
    )

    # Title — Chelsea Market, big, centered
    requests += _create_text_box(
        slide_id, _gen_id("ttl", idx, 0),
        x_in=0.35, y_in=1.30, w_in=6.8, h_in=1.40,
        text=cert["title"], font_size_pt=36,
        font_family=TITLE_FONT, align="CENTER",
    )
    # Recipient label
    requests += _create_text_box(
        slide_id, _gen_id("rec", idx, 0),
        x_in=0.5, y_in=3.10, w_in=6.5, h_in=0.50,
        text=cert["recipient_field_label"],
        font_size_pt=18, font_family=BODY_FONT, align="CENTER",
    )
    requests += _create_underline(
        slide_id, _gen_id("rec_ln", idx, 0),
        x_in=1.50, y_in=3.95, w_in=4.50, weight_pt=1.5,
    )

    # Achievement_text — dynamically sized. G3's achievement_text is ~600
    # chars and used to overflow into the "By completing…" header below.
    # Truncate to a known per-line budget AND shrink font when long, so
    # the box height fits within the budget between recipient_underline
    # (y=3.95) and the skills-header (which we now place dynamically).
    ACH_Y = 4.30
    ACH_MAX_H = 1.60          # max budget for achievement text
    ACH_LINE_IN = 0.28        # 16pt line height
    ach_w_chars = 70          # ~16pt across 6.5"
    # Pick a font that lets the text fit inside ACH_MAX_H without overflow.
    ach_text = cert["achievement_text"]
    ach_pt = 16
    if _est_lines(ach_text, ach_w_chars) * ACH_LINE_IN > ACH_MAX_H:
        # Try 14pt (line height ~0.24")
        ach_pt = 14
        ACH_LINE_IN = 0.24
        ach_w_chars = 80
        if _est_lines(ach_text, ach_w_chars) * ACH_LINE_IN > ACH_MAX_H:
            # Last resort — 12pt
            ach_pt = 12
            ACH_LINE_IN = 0.20
            ach_w_chars = 95
            # And truncate in case still long
            ach_text = _truncate_to_height(ach_text, ACH_MAX_H, ACH_LINE_IN, ach_w_chars)

    n_lines = _est_lines(ach_text, ach_w_chars)
    ach_h = min(ACH_MAX_H, max(0.6, n_lines * ACH_LINE_IN + 0.10))
    requests += _create_text_box(
        slide_id, _gen_id("ach", idx, 0),
        x_in=0.5, y_in=ACH_Y, w_in=6.5, h_in=ach_h,
        text=ach_text,
        font_size_pt=ach_pt, font_family=BODY_FONT, align="CENTER",
    )

    # Skills header & body now flow from the bottom of the achievement box
    # rather than fixed positions. Keeps a 0.20" gap above the header.
    sk_hdr_y = ACH_Y + ach_h + 0.20
    requests += _create_text_box(
        slide_id, _gen_id("sk_hdr", idx, 0),
        x_in=0.75, y_in=sk_hdr_y, w_in=6.0, h_in=0.45,
        text="By completing this unit, this student has demonstrated:",
        font_size_pt=14, font_family=BODY_FONT, bold=True,
    )
    skills_text = "\n".join(f"  •  {s}" for s in cert["skills_demonstrated"])
    skills_y = sk_hdr_y + 0.55
    # Cap skills box height so the Coco corner image (y=8.15) doesn't overlap.
    skills_h = max(0.80, 8.10 - skills_y)
    # Truncate skills text if it would overflow
    skills_text = _truncate_to_height(skills_text, skills_h, line_in=0.28, w_chars=70)
    requests += _create_text_box(
        slide_id, _gen_id("skills", idx, 0),
        x_in=0.85, y_in=skills_y, w_in=5.5, h_in=skills_h,
        text=skills_text, font_size_pt=14, font_family=BODY_FONT,
    )
    # Coco — bottom-right corner
    coco_path = composed_dir / "AS_CERT_COCO.png"
    if coco_path.exists():
        cid = uploader.upload(coco_path, "as_cert_coco.png")
        requests.append(_create_image(slide_id, _gen_id("coco", idx, 0), cid,
                                      x_in=5.50, y_in=8.15, w_in=1.40, h_in=1.20))
    return requests


def build_marketplace_slide(idx: int, slide_id: str, mk: dict) -> list[dict]:
    requests = [_create_slide_request(slide_id)]
    # Title (Lexend bold)
    requests += _create_text_box(
        slide_id, _gen_id("ttl", idx, 0),
        x_in=0.35, y_in=0.30, w_in=6.8, h_in=0.65,
        text="Marketplace Listing", font_size_pt=24, bold=True, font_family=HEADING_FONT,
    )
    desc_short = _first_sentence(mk["short_description"], max_chars=320)
    desc_text = f"{desc_short}"
    requests += _create_text_box(
        slide_id, _gen_id("desc", idx, 0),
        x_in=0.35, y_in=1.10, w_in=6.8, h_in=1.50,
        text=desc_text, font_size_pt=15, font_family=BODY_FONT,
    )
    meta_text = (
        f"Price: ${mk['suggested_price_cad']} CAD\n"
        f"Pages: ~{mk['pages_total_estimate']}\n"
        f"Classroom time: {mk['classroom_time_total_minutes']} min\n"
        f"Teacher prep: ~{mk['teacher_prep_time_minutes']} min"
    )
    requests += _create_text_box(
        slide_id, _gen_id("meta", idx, 0),
        x_in=0.35, y_in=2.75, w_in=6.8, h_in=1.20,
        text=meta_text, font_size_pt=14, font_family=BODY_FONT,
    )
    tags_text = f"Top tags: {', '.join(mk['tags'][:8])}"
    requests += _create_text_box(
        slide_id, _gen_id("tags", idx, 0),
        x_in=0.35, y_in=4.05, w_in=6.8, h_in=0.55,
        text=tags_text, font_size_pt=13, font_family=BODY_FONT, color=GRAY_LABEL,
    )
    requests += _create_text_box(
        slide_id, _gen_id("inc_hdr", idx, 0),
        x_in=0.35, y_in=4.75, w_in=6.8, h_in=0.45,
        text="What's included:", font_size_pt=16, bold=True, font_family=HEADING_FONT,
    )
    inc_items = [_first_sentence(x, max_chars=160) for x in mk["what_is_included"][:12]]
    inc_text = "\n".join(f"   •  {x}" for x in inc_items)
    inc_box_h = 4.40
    inc_box_w = 6.8
    inc_pt = _fit_font_size(inc_text, base_pt=13, w_in=inc_box_w, h_in=inc_box_h, min_pt=10)
    requests += _create_text_box(
        slide_id, _gen_id("inc", idx, 0),
        x_in=0.35, y_in=5.30, w_in=inc_box_w, h_in=inc_box_h,
        text=inc_text, font_size_pt=inc_pt, font_family=BODY_FONT,
    )
    return requests


# ── Top-level driver ──────────────────────────────────────────────────────

def render_validation_pages(unit_dir: Path, dpi: int = 150,
                             out_dir: Path | None = None) -> list[Path]:
    """Render the unit's validation_export.pdf to one PNG per slide so that
    the runner (Claude in chat) can visually inspect every slide before
    declaring the deck shipped.

    The text-based pre-flight (`validate_unit_for_slides`) catches signal-
    level problems but cannot see actual rendered overlap, truncation, or
    image-placeholder failures. The fix-iterate-fix loop on G3 exposed
    five separate visual issues (lesson Exit cuts, worksheet title
    overlap, hero placeholder PNGs, certificate text bleed, manipulative
    prep overflow) that text heuristics alone never flagged. Visual
    inspection is therefore a mandatory step in the unit-generation
    workflow — see CLAUDE.md "Final-output protocol" for the runbook.

    Returns the list of generated PNG paths (one per slide), in slide
    order. Each is suitable for the Read tool.
    """
    import shutil, subprocess
    pdf_path = unit_dir / "validation_export.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(
            f"validation_export.pdf missing at {pdf_path}. "
            f"Run build_unit_deck() first — it writes the PDF as a side effect."
        )
    if out_dir is None:
        # Use a temp dir under the unit folder so the renderings don't
        # accumulate next to the deck.
        out_dir = unit_dir / "_inspect_pages"
    if out_dir.exists():
        shutil.rmtree(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    if shutil.which("pdftoppm") is None:
        raise RuntimeError(
            "pdftoppm not installed — install poppler "
            "(brew install poppler / apt install poppler-utils)"
        )
    subprocess.run(
        ["pdftoppm", "-png", "-r", str(dpi), str(pdf_path), str(out_dir / "slide")],
        check=True,
    )
    return sorted(out_dir.glob("slide-*.png"))


# Per-slide-type visual-inspection checklist. The runner walks through this
# while viewing each rendered PNG; failures should trigger a layout fix or a
# content edit before declaring the deck shipped. Maintain alongside the
# fixes we've already added — anything you discover visually that the text
# pre-flight missed should land here.
VISUAL_INSPECTION_CHECKLIST: dict[str, list[str]] = {
    "cover": [
        "Grade label matches blueprint.grade (NOT hardcoded 'Kindergarten')",
        "Subtitle box does not strikethrough or wrap awkwardly",
        "Coco character at the bottom is not clipped",
    ],
    "overview": [
        "Title fits without overlapping the Grade/Strand line below",
        "unit_learning_goal does not push the lesson-arc table off the slide",
        "All 5 lesson-arc rows visible",
    ],
    "lesson": [
        "Title fits the title bar without wrapping into metadata",
        "Metadata block (Goal/Expectation/Learning Goal/Materials) visible in full",
        "All action bullets visible (not cut at body box edge)",
        "All consolidation prompts visible",
        "Exit text ends with a complete sentence or '…' truncation marker — NOT mid-word at the page edge",
        "Sections do not overlap each other vertically",
    ],
    "worksheet": [
        "Title sits ENTIRELY above the Learning Goal box (no second line bleeding into it)",
        "Hero image is real artwork, not a grey 'placeholder' PNG showing the asset_id",
        "Footer '(N parts on the printable worksheet)' visible at bottom",
    ],
    "manipulative": [
        "Title + asset_id visible at top",
        "Image fits and is not cut by the spec/prep columns below",
        "TEACHER PREP column ends with a complete step or '…' marker",
    ],
    "rubric": [
        "Table fits within slide; no row truncated",
        "Level 1-4 descriptors all readable",
    ],
    "certificate": [
        "Title 'Junior … Certificate' on one or two lines, fully visible",
        "achievement_text body does NOT touch or overlap 'By completing this unit…'",
        "Skills bullets do not overlap the Coco corner image",
    ],
    "marketplace": [
        "What's-included list visible to the bottom margin",
        "Price/pages/time meta block readable",
    ],
    "terms": [
        "Standard ToU text fully visible",
    ],
}


def validate_unit_for_slides(unit_dir: Path) -> list[str]:
    """Pre-flight content-size check. Returns a list of warnings about content
    that is likely to overflow or compress poorly on the slide layout.

    Run this before `build_unit_deck()` to catch problems Pydantic schemas
    don't see — e.g. teacher scripts that exceed the slide budget for a given
    grade. Empty list = clean.

    Added 2026-04-28 after the G1 gap-finding run surfaced cover/overview
    overflows. Treats warnings as non-blocking; the build still proceeds, but
    the next session can decide whether to revisit content."""
    warnings: list[str] = []

    bp_path = unit_dir / "0_blueprint.json"
    if not bp_path.exists():
        return ["blueprint not yet generated"]
    bp = json.loads(bp_path.read_text(encoding="utf-8"))

    # 1. Cover-slide subtitle length (warn if > 90 chars — auto-fit caps at 16pt)
    if len(bp.get("thematic_title", "")) > 90:
        warnings.append(
            f"cover: thematic_title is {len(bp['thematic_title'])} chars; subtitle font will shrink to ~14pt and may still wrap awkwardly")
    # 2. Overview title (Unit Overview: <title>) length
    ov_title = f"Unit Overview: {bp.get('thematic_title', '')}"
    if len(ov_title) > 110:
        warnings.append(
            f"overview: title is {len(ov_title)} chars; will shrink to ~18pt and push the lesson-arc table down")
    # 3. unit_learning_goal length
    goal = bp.get("unit_learning_goal", "")
    if len(goal) > 320:
        warnings.append(
            f"overview: unit_learning_goal is {len(goal)} chars; will consume >2 lines and compress the lesson-arc table")

    # 4. Lessons — teacher script length and consolidation prompt count
    for lp_path in sorted(unit_dir.glob("1_lesson_*.json")):
        lp = json.loads(lp_path.read_text(encoding="utf-8"))
        n = lp.get("lesson_number", "?")
        mo = lp.get("minds_on", {})
        script = mo.get("teacher_script", "")
        # First-paragraph length (used on slide). Warn over ~700 chars even
        # though the slide truncates at 360 — long scripts often hide content
        # the teacher wants visible.
        para = script.split("\n\n")[0] if "\n\n" in script else script
        if len(para) > 700:
            warnings.append(
                f"lesson_{n:02d}: minds_on first paragraph is {len(para)} chars; slide will truncate at 360 — consider splitting into two paragraphs")
        action = lp.get("action", {})
        steps = action.get("steps", [])
        if len(steps) > 6:
            warnings.append(
                f"lesson_{n:02d}: action has {len(steps)} steps; lesson slide will compress them to fit")
        cons = lp.get("consolidation", {})
        prompts = cons.get("discussion_prompts", [])
        if len(prompts) > 5:
            warnings.append(
                f"lesson_{n:02d}: consolidation has {len(prompts)} prompts; slide will compress to fit")

    # 5. Marketplace — what_is_included length
    mk_path = unit_dir / "6_marketplace.json"
    if mk_path.exists():
        mk = json.loads(mk_path.read_text(encoding="utf-8"))
        inc = mk.get("what_is_included", [])
        if len(inc) > 14:
            warnings.append(
                f"marketplace: {len(inc)} 'what_is_included' items; slide will only render first 12")

    # 6. Reflection sheet — Prompt 2 yes/no row count (slide can hold ~6 rows comfortably)
    fr_path = unit_dir / "4_formative_reflection.json"
    if fr_path.exists():
        fr = json.loads(fr_path.read_text(encoding="utf-8"))
        for prompt in fr.get("reflection_sheet", {}).get("prompts", []):
            opts = prompt.get("options") or []
            if prompt.get("response_type") == "circle_yes_no" and len(opts) > 6:
                warnings.append(
                    f"reflection_sheet.prompt_{prompt.get('prompt_number')}: {len(opts)} yes/no rows; consider splitting into two prompts")

    # 7. Worksheets — title length (Chelsea Market wraps a second line into the
    # Learning Goal box at >25 chars unless auto-fit shrinks the font).
    # Emit a warning at 30+ chars so the runner notices the title font will
    # drop below 24pt. Also warn if the chosen "hero" composite is suspiciously
    # small (likely a placeholder PNG, not real artwork).
    composed_dir = unit_dir / "composed"
    PLACEHOLDER_BYTES = 2500
    for ws_path in sorted(unit_dir.glob("2_worksheet_*.json")):
        ws = json.loads(ws_path.read_text(encoding="utf-8"))
        wn = ws.get("lesson_number", "?")
        title = ws.get("worksheet_title", "")
        if len(title) > 30:
            warnings.append(
                f"worksheet_{wn:02d}: title is {len(title)} chars; will shrink "
                f"below 24pt to fit single-line — consider a shorter title for visual weight")
        # Hero candidate scan: every part image_placeholder file size
        if composed_dir.exists():
            placeholder_hits: list[str] = []
            for part in ws.get("pages", [{}])[0].get("parts", []):
                for ph in part.get("image_placeholders", []) or []:
                    p = composed_dir / f"{ph['id']}.png"
                    if p.exists() and p.stat().st_size < PLACEHOLDER_BYTES:
                        placeholder_hits.append(ph["id"])
            if placeholder_hits and len(placeholder_hits) == sum(
                1 for part in ws.get("pages", [{}])[0].get("parts", [])
                for ph in part.get("image_placeholders", []) or []
            ):
                # ALL composites for this worksheet are placeholders
                warnings.append(
                    f"worksheet_{wn:02d}: every composite under {PLACEHOLDER_BYTES}B "
                    f"({placeholder_hits}); slide will fall back to placeholder graphics — "
                    f"likely missing real artwork")

    # 8. Certificate — achievement_text length (G3's was 700+ chars and overflowed
    # into the skills header before the dynamic-sizing fix). Warn at 600+.
    asu_path = unit_dir / "5_assessment_suite.json"
    if asu_path.exists():
        asu = json.loads(asu_path.read_text(encoding="utf-8"))
        cert = asu.get("certificate", {})
        ach = cert.get("achievement_text", "") or ""
        if len(ach) > 600:
            warnings.append(
                f"certificate: achievement_text is {len(ach)} chars; will shrink "
                f"below 14pt and risk pushing the skills header into the skills body")
        skills = cert.get("skills_demonstrated", []) or []
        if len(skills) > 6:
            warnings.append(
                f"certificate: {len(skills)} skills_demonstrated bullets; "
                f"slide caps at ~6 before the Coco corner image overlap")

    # 9. Manipulatives — teacher_prep_steps total length (slide budget 3.70" tall)
    mn_path = unit_dir / "3_manipulatives.json"
    if mn_path.exists():
        mn = json.loads(mn_path.read_text(encoding="utf-8"))
        for asset in mn.get("manipulatives", []) or []:
            steps = asset.get("teacher_prep_steps", []) or []
            total_chars = sum(len(s) for s in steps[:5])
            if total_chars > 700:
                warnings.append(
                    f"manipulative {asset.get('asset_id','?')}: teacher_prep_steps "
                    f"total {total_chars} chars; slide will truncate to fit")

    return warnings


def build_unit_deck(unit_dir: Path, run_preflight: bool = True,
                    drive_parent_folder_id: str | None = None) -> str:
    """Build a Slides deck for the unit. Returns the deck URL.

    The deck and any Drive-uploaded composite images are placed in a
    per-unit subfolder under ``drive_parent_folder_id`` (defaults to
    ``UNIT_DECK_PARENT_FOLDER_ID`` — the TCE shared folder). Pass
    ``drive_parent_folder_id=None`` and set the env var to "" to fall back
    to Drive root (legacy behaviour for ad-hoc one-off builds).

    `run_preflight=True` (default) calls `validate_unit_for_slides()` before
    building and prints warnings to stdout. Warnings are non-blocking — the
    build proceeds either way. Pass `run_preflight=False` to skip.
    """
    if run_preflight:
        warns = validate_unit_for_slides(unit_dir)
        if warns:
            print(f"Pre-flight slide-layout warnings ({len(warns)}):")
            for w in warns:
                print(f"  ⚠ {w}")
        else:
            print("Pre-flight slide-layout check: ✓ clean")
    creds = get_credentials()
    slides = build("slides", "v1", credentials=creds)
    drive = build("drive", "v3", credentials=creds)

    # Load blueprint early — needed both for the per-unit subfolder name and
    # for slide content downstream.
    bp = json.loads((unit_dir / "0_blueprint.json").read_text())

    # ── Publication gate (rubric grade) ─────────────────────────────────
    # The deck only lands in the public TCE folder if the unit has a
    # passing RubricGrade (≥ 17/20). Failed or ungraded units route to a
    # `_drafts` sibling folder so the deck is still inspectable for the
    # Appearance score, but unmistakably not "shipped".
    grade_path = unit_dir / "7_rubric_grade.json"
    grade_status: str
    grade_overall: int | None = None
    if grade_path.exists():
        try:
            _g = json.loads(grade_path.read_text())
            grade_status = _g.get("status", "unknown")
            grade_overall = _g.get("overall_score")
        except Exception:
            grade_status = "unreadable"
    else:
        grade_status = "missing"
    is_published = (grade_status == "pass")
    print(f"Rubric grade: status={grade_status!r} overall={grade_overall!r} → "
          f"{'PUBLISH to shared folder' if is_published else 'route to _drafts subfolder'}")

    # Resolve where the deck + assets land. A per-unit subfolder under the
    # configured parent keeps shared-drive listings tidy and matches the
    # pattern already in use ("Test - Pattern Unit", "Coding Unit", …).
    parent_id = drive_parent_folder_id if drive_parent_folder_id is not None \
        else (UNIT_DECK_PARENT_FOLDER_ID or None)
    unit_folder_id: str | None = None
    if parent_id:
        unit_folder_name = bp.get("thematic_title") or bp.get("unit_id") or unit_dir.name
        if not is_published:
            # Failed/ungraded → put both the per-unit subfolder under "_drafts"
            # so it's siloed from shipped content.
            drafts_id = _find_or_create_subfolder(drive, parent_id, "_drafts")
            unit_folder_id = _find_or_create_subfolder(drive, drafts_id, unit_folder_name)
            print(f"Drive destination: parent={parent_id} → _drafts → '{unit_folder_name}' (id={unit_folder_id})")
        else:
            unit_folder_id = _find_or_create_subfolder(drive, parent_id, unit_folder_name)
            print(f"Drive destination: parent={parent_id} → unit_folder='{unit_folder_name}' (id={unit_folder_id})")
    else:
        print("Drive destination: root (no parent folder configured)")

    uploader = DriveUploader(drive, folder_id=unit_folder_id)

    composed_dir = unit_dir / "composed"
    if not composed_dir.exists() or not any(composed_dir.iterdir()):
        from pipeline.compose import compose_for_unit
        compose_for_unit(unit_dir)

    # Remaining stage data
    mk = json.loads((unit_dir / "6_marketplace.json").read_text())
    lessons = [json.loads(p.read_text()) for p in sorted(unit_dir.glob("1_lesson_*.json"))]
    worksheets = [json.loads(p.read_text()) for p in sorted(unit_dir.glob("2_worksheet_*.json"))]
    manip = json.loads((unit_dir / "3_manipulatives.json").read_text())
    fr = json.loads((unit_dir / "4_formative_reflection.json").read_text())
    asu = json.loads((unit_dir / "5_assessment_suite.json").read_text())

    # Create the presentation in PORTRAIT 7.5×10.
    # The Slides API ignores pageSize at create time, so we use a workaround:
    # 1. Generate a blank portrait .pptx via python-pptx
    # 2. Upload to Drive and convert to a Google Slides presentation
    # 3. The resulting deck has the portrait dimensions
    from datetime import datetime
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Emu
    import io

    stamp = datetime.now().strftime("%Y-%m-%d %H%M")
    title = f"{bp['thematic_title']} — Unit Deck ({stamp})"
    print(f"Creating presentation: {title}")

    # Build a blank portrait .pptx in memory
    pptx = PPTXPresentation()
    pptx.slide_width = Emu(SLIDE_WIDTH_EMU)
    pptx.slide_height = Emu(SLIDE_HEIGHT_EMU)
    blank_layout = pptx.slide_layouts[6]
    pptx.slides.add_slide(blank_layout)
    buf = io.BytesIO()
    pptx.save(buf)
    buf.seek(0)

    # Upload to Drive as a Google Slides (server-side conversion via mimeType)
    from googleapiclient.http import MediaIoBaseUpload
    media = MediaIoBaseUpload(
        buf, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=False,
    )
    deck_body: dict = {
        "name": title,
        "mimeType": "application/vnd.google-apps.presentation",  # convert on upload
    }
    if unit_folder_id:
        deck_body["parents"] = [unit_folder_id]
    drive_file = drive.files().create(
        body=deck_body,
        media_body=media,
        fields="id",
    ).execute()
    pres_id = drive_file["id"]
    pres = slides.presentations().get(presentationId=pres_id).execute()

    # Verify portrait
    ps = pres.get("pageSize", {})
    print(f"  pageSize: {ps['width']['magnitude']/914400:.1f} × {ps['height']['magnitude']/914400:.1f} in")

    # Delete the default first slide
    if pres.get("slides"):
        delete_default = [{"deleteObject": {"objectId": pres["slides"][0]["objectId"]}}]
        slides.presentations().batchUpdate(presentationId=pres_id, body={"requests": delete_default}).execute()

    # Build all the slide creation + content requests
    requests: list[dict] = []
    table_fills: list[tuple] = []  # deferred: (table_id, fill_callable)

    idx = 0

    # 1. Cover
    cover_id = f"slide_cover"
    requests += build_cover_slide(idx, cover_id, bp, mk, uploader, composed_dir)
    idx += 1

    # 2. Overview
    overview_id = f"slide_overview"
    ov_reqs, ov_table = build_overview_slide(idx, overview_id, bp)
    requests += ov_reqs

    # Defer table cell-fill until after table creation
    def fill_overview(table_id):
        fills = []
        for col, label in enumerate(["Lesson", "Lesson Title", "Learning Goal"]):
            fills += _insert_table_text(table_id, 0, col, label, font_size_pt=12, bold=True)
        for r, entry in enumerate(bp["lesson_arc"], 1):
            fills += _insert_table_text(table_id, r, 0, f"Day {entry['lesson_number']}", font_size_pt=12, bold=True)
            fills += _insert_table_text(table_id, r, 1, entry["lesson_title"], font_size_pt=12, bold=True)
            fills += _insert_table_text(table_id, r, 2, entry["student_learning_goal"], font_size_pt=12)
        return fills
    table_fills.append((ov_table, fill_overview))
    idx += 1

    # 3. Lesson plans (5 — or 10 for G3+, since dense lessons split onto 2 slides)
    grade_for_split = bp.get("grade")
    for lp in lessons:
        sid_base = f"slide_lesson_{lp['lesson_number']:02d}"
        reqs, n_slides = build_lesson_slides(idx, sid_base, lp, grade=grade_for_split)
        requests += reqs
        idx += n_slides

    # 4. Worksheets (5)
    for ws in worksheets:
        sid = f"slide_worksheet_{ws['lesson_number']:02d}"
        requests += build_worksheet_slide(idx, sid, ws, uploader, composed_dir)
        idx += 1

    # 5. Manipulatives (one slide each)
    for asset in manip["assets"]:
        sid = f"slide_manip_{asset['asset_id']}"
        requests += build_manipulative_slide(idx, sid, asset, uploader, composed_dir)
        idx += 1

    # 6. Rubric
    rubric_sid = "slide_rubric"
    rub_reqs, rub_table = build_rubric_slide(idx, rubric_sid, asu["summative_rubric"], bp)
    requests += rub_reqs

    def fill_rubric(table_id):
        fills = []
        rub = asu["summative_rubric"]
        grade = bp.get("grade", "")
        # Headers — bigger, bolder
        headers = ["Expectation", "Level 1\n(Beginning)", "Level 2\n(Developing)", "Level 3\n(Achieving)", "Level 4\n(Exceeding)"]
        for col, label in enumerate(headers):
            fills += _insert_table_text(table_id, 0, col, label, font_size_pt=11, bold=True)

        # Lookup tables for short on-slide descriptors. Keys can be either:
        #   (code, level)        — grade-blind default
        #   (grade, code, level) — grade-specific override (preferred)
        # Lookup precedence: (grade, code, level) → (code, level) → first-sentence fallback.
        # SHORT_EXP follows the same pattern: (code) and (grade, code).
        # Refactored 2026-04-29 to add per-grade keys; G1 and G2 used to share
        # the same C1.x descriptors which lost grade specificity.

        SHORT_EXP = {
            # ── Generic per-code defaults (any grade if no override exists) ──
            # Kindergarten Patterns (A7.x)
            "A7.1": "Identify, describe, extend, create patterns",
            "A7.2": "Identify the core; explain why it matters",
            "A7.3": "Show same core in different ways",
            "A7.4": "Extend in both directions; find missing",
            # Generic Algebra Patterns (C1.x)
            "C1.1": "Identify, describe, extend repeating patterns",
            "C1.2": "Translate patterns across representations",
            "C1.3": "Determine rules; justify predictions; find missing",
            "C1.4": "Use patterns to describe whole-number relationships",

            # ── Grade-specific overrides ──
            # Grade 1
            ("Grade 1", "C1.1"): "Identify, describe, extend repeating patterns",
            ("Grade 1", "C1.2"): "Translate patterns (objects, sounds, movements, pictures)",
            ("Grade 1", "C1.3"): "Determine rules (AB/ABB/AAB/ABC); justify; find missing",
            ("Grade 1", "C1.4"): "Use patterns to describe whole numbers up to 50",
            # Grade 2
            ("Grade 2", "C1.1"): "Identify repeating, growing, shrinking patterns",
            ("Grade 2", "C1.2"): "Translate across reps (shapes, numbers, tables of values)",
            ("Grade 2", "C1.3"): "Determine rules (+N/-N); justify; find missing in any type",
            ("Grade 2", "C1.4"): "Use patterns to describe whole numbers up to 100",
            # Grade 3 (forward-looking; ready for the upcoming G3 unit)
            ("Grade 3", "C1.1"): "Identify repeating elements and operations in patterns",
            ("Grade 3", "C1.2"): "Translate patterns; build tables of values",
            ("Grade 3", "C1.3"): "Determine rules (incl. growing/shrinking +/-); justify",
            ("Grade 3", "C1.4"): "Use patterns to describe whole numbers up to 1000",
        }

        SHORT_LEVEL = {
            # ── Kindergarten (A7.x — no grade override needed; only K uses A7.x) ──
            ("A7.1", 1): "Names what's next only with prompts.",
            ("A7.1", 2): "Predicts most of the time. Creates AB with help.",
            ("A7.1", 3): "Confidently predicts and creates an AB pattern.",
            ("A7.1", 4): "Creates 3-element cores; spots patterns in environment.",
            ("A7.2", 1): "Cannot consistently isolate the core.",
            ("A7.2", 2): "Identifies AB cores with adult support.",
            ("A7.2", 3): "Identifies AB, ABB, ABC cores; explains repetition.",
            ("A7.2", 4): "Articulates why the core matters; predicts beyond.",
            ("A7.3", 1): "Treats different reps as different patterns.",
            ("A7.3", 2): "Produces a 2nd version with help; structure may drift.",
            ("A7.3", 3): "Translates one core into 2 different representations.",
            ("A7.3", 4): "Translates into 3+ reps; articulates core vs surface.",
            ("A7.4", 1): "Cannot fill missing elements.",
            ("A7.4", 2): "Fills missing at the end; struggles with start.",
            ("A7.4", 3): "Fills missing anywhere; extends forward correctly.",
            ("A7.4", 4): "Extends both directions; explains what came before.",

            # ── Generic C1.x defaults (used when no grade override exists) ──
            ("C1.1", 1): "Names what's next only with prompts.",
            ("C1.1", 2): "Predicts most of the time. Creates 2-element pattern with help.",
            ("C1.1", 3): "Confidently predicts and creates 2- or 3-element patterns.",
            ("C1.1", 4): "Creates 3- or 4-element cores; spots patterns across reps.",
            ("C1.2", 1): "Treats different reps as different patterns (surface only).",
            ("C1.2", 2): "Produces one 2nd rep with help; structure may drift.",
            ("C1.2", 3): "Independently translates a rule into 2 different reps.",
            ("C1.2", 4): "Translates into 3+ reps including invented ones.",
            ("C1.3", 1): "Cannot consistently name a rule or fill missing.",
            ("C1.3", 2): "Names a rule with prompts; fills missing at end only.",
            ("C1.3", 3): "Names rule independently; fills missing anywhere; justifies.",
            ("C1.3", 4): "Names invented rules; extends both directions independently.",
            ("C1.4", 1): "Does not yet translate animals to numbers.",
            ("C1.4", 2): "Builds 2-element number parade with small numbers, with help.",
            ("C1.4", 3): "Builds repeating number parade; explains relationship.",
            ("C1.4", 4): "Uses larger numbers; predicts cells far beyond.",

            # ── Grade 1 overrides ──
            ("Grade 1", "C1.1", 1): "Names what's next only with prompts.",
            ("Grade 1", "C1.1", 2): "Predicts most of the time. Creates 2-element with help.",
            ("Grade 1", "C1.1", 3): "Confidently predicts and creates 2- or 3-element patterns.",
            ("Grade 1", "C1.1", 4): "Creates 3- or 4-element cores; spots patterns across reps.",
            ("Grade 1", "C1.2", 1): "Treats different reps as different patterns.",
            ("Grade 1", "C1.2", 2): "Produces one 2nd rep with help; structure may drift.",
            ("Grade 1", "C1.2", 3): "Translates rule (AB/ABB/AAB/ABC) into 2 different reps.",
            ("Grade 1", "C1.2", 4): "Translates into 3+ reps including invented ones.",
            ("Grade 1", "C1.3", 1): "Cannot consistently name a rule or fill missing.",
            ("Grade 1", "C1.3", 2): "Names AB with prompts; fills missing at end only.",
            ("Grade 1", "C1.3", 3): "Names AB/ABB/AAB/ABC; fills missing anywhere; justifies.",
            ("Grade 1", "C1.3", 4): "Names invented rules; extends both directions independently.",
            ("Grade 1", "C1.4", 1): "Does not yet translate animals to numbers.",
            ("Grade 1", "C1.4", 2): "Builds 2-element number parade with small numbers, with help.",
            ("Grade 1", "C1.4", 3): "Builds repeating number parade up to 50; explains relationship.",
            ("Grade 1", "C1.4", 4): "Uses larger numbers within 50; predicts cells far beyond.",

            # ── Grade 2 overrides (growing/shrinking + table of values) ──
            ("Grade 2", "C1.1", 1): "Identifies pattern type only with prompts.",
            ("Grade 2", "C1.1", 2): "Identifies the 3 types most of the time; informal language.",
            ("Grade 2", "C1.1", 3): "Confidently identifies repeating, growing, shrinking; spots in real life.",
            ("Grade 2", "C1.1", 4): "Identifies pattern types in unfamiliar contexts; combined types.",
            ("Grade 2", "C1.2", 1): "Treats reps as different patterns; no table of values.",
            ("Grade 2", "C1.2", 2): "Translates with help; table may miss headers or mismatch.",
            ("Grade 2", "C1.2", 3): "Translates across parade, numbers, and table of values.",
            ("Grade 2", "C1.2", 4): "Translates into invented reps; uses table to predict relationships.",
            ("Grade 2", "C1.3", 1): "Cannot name +N / -N rule; fills missing randomly.",
            ("Grade 2", "C1.3", 2): "Names +N / -N with prompts; fills missing at end only.",
            ("Grade 2", "C1.3", 3): "Names rule, fills missing anywhere, predicts term 10 by skip-counting.",
            ("Grade 2", "C1.3", 4): "Predicts term 20+ efficiently; recognizes 100-ceiling constraints.",
            ("Grade 2", "C1.4", 1): "Does not yet build a number parade up to 100.",
            ("Grade 2", "C1.4", 2): "Builds with small numbers (1-20); relationship unclear.",
            ("Grade 2", "C1.4", 3): "Builds number parade up to 100; plans starting number; explains relationship.",
            ("Grade 2", "C1.4", 4): "Uses larger rules; recognizes factors of 100; generalizes term-value relationship.",
        }

        def lookup_short_exp(code: str) -> str:
            return (
                SHORT_EXP.get((grade, code))
                or SHORT_EXP.get(code)
                or _first_sentence(
                    next((row_["expectation_text"] for row_ in rub["rows"]
                          if row_["expectation_code"] == code), code),
                    max_chars=50,
                )
            )

        def lookup_short_level(code: str, level: int, level_descriptor: str) -> str:
            return (
                SHORT_LEVEL.get((grade, code, level))
                or SHORT_LEVEL.get((code, level))
                or _first_sentence(level_descriptor, max_chars=80)
            )

        for r, row in enumerate(rub["rows"], 1):
            exp_short = lookup_short_exp(row["expectation_code"])
            fills += _insert_table_text(table_id, r, 0,
                                        f"{row['expectation_code']}\n{exp_short}",
                                        font_size_pt=10, bold=True)
            for col_n in range(1, 5):
                lvl_desc = next(l["descriptor"] for l in row["levels"]
                                if l["level_number"] == col_n)
                desc = lookup_short_level(row["expectation_code"], col_n, lvl_desc)
                fills += _insert_table_text(table_id, r, col_n, desc, font_size_pt=10)
        return fills
    table_fills.append((rub_table, fill_rubric))
    idx += 1

    # 7. Certificate
    requests += build_certificate_slide(idx, "slide_certificate",
                                        asu["certificate"], bp, uploader, composed_dir)
    idx += 1

    # 8. Marketplace listing
    requests += build_marketplace_slide(idx, "slide_marketplace", mk)
    idx += 1

    # Apply slide-creation requests in batches (Slides API limits batch sizes)
    print(f"Sending {len(requests)} slide-creation requests for {idx} slides...")
    BATCH = 100
    for i in range(0, len(requests), BATCH):
        slides.presentations().batchUpdate(
            presentationId=pres_id,
            body={"requests": requests[i:i + BATCH]}
        ).execute()

    # Now apply deferred table fills
    fill_requests = []
    for table_id, fill_fn in table_fills:
        fill_requests += fill_fn(table_id)
    if fill_requests:
        print(f"Filling {len(fill_requests)} table cell requests...")
        for i in range(0, len(fill_requests), BATCH):
            slides.presentations().batchUpdate(
                presentationId=pres_id,
                body={"requests": fill_requests[i:i + BATCH]}
            ).execute()

    url = f"https://docs.google.com/presentation/d/{pres_id}/edit"
    print(f"\n✓ Deck created: {url}")

    # Export to PDF for visual validation
    pdf_path = unit_dir / "validation_export.pdf"
    print(f"Exporting deck to PDF for validation: {pdf_path}")
    pdf_content = drive.files().export(
        fileId=pres_id, mimeType="application/pdf"
    ).execute()
    pdf_path.write_bytes(pdf_content)
    print(f"  ✓ PDF saved ({len(pdf_content):,} bytes)")

    # Clean up composite PNGs from the unit's Drive folder. Slides has already
    # downloaded server-side copies of every embedded image at insert time, so
    # the source files in Drive are no longer needed. Deleting them leaves
    # ONLY the .gslides deck visible in the per-unit folder — the buyer/teacher-
    # facing artifact. Skipped when no folder is configured (legacy local-build
    # path) and skipped for failed/ungraded decks routed to _drafts (which
    # users may still want to inspect with the source images intact).
    if unit_folder_id and is_published:
        print(f"Cleaning up composite assets from unit folder…")
        try:
            asset_query = (
                f"'{unit_folder_id}' in parents and "
                f"mimeType != 'application/vnd.google-apps.presentation' and "
                f"mimeType != 'application/vnd.google-apps.folder' and "
                f"trashed = false"
            )
            assets_resp = drive.files().list(
                q=asset_query,
                fields="files(id,name,mimeType)",
                pageSize=200,
            ).execute()
            assets = assets_resp.get("files", [])
            n_deleted = 0
            for f in assets:
                try:
                    drive.files().delete(fileId=f["id"]).execute()
                    n_deleted += 1
                except Exception as e:
                    print(f"    (could not delete {f['name']!r}: {e})")
            print(f"  ✓ Removed {n_deleted}/{len(assets)} composite asset(s) "
                  f"from the unit folder; deck remains.")
        except Exception as e:
            # Don't fail the build over cleanup; just log it.
            print(f"  ⚠ Asset cleanup failed (non-blocking): {e}")

    return url


if __name__ == "__main__":
    import sys
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else \
        PROJECT_ROOT / "generated_units/batch_1/k_patterns_pattern_parade"
    print(build_unit_deck(target))
